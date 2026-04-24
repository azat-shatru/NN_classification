"""
PPT export — native python-pptx charts for Stage 2.6 Charts Portal.
Builds actual chart objects (not images) so they remain editable in PowerPoint.
"""
import io
import re
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
_W       = Inches(13.33)
_H       = Inches(7.5)
_MARGIN  = Inches(0.3)
_GAP     = Inches(0.2)
_TITLE_H = Inches(0.65)
_TOP     = Inches(0.9)

# ── Colour palettes ───────────────────────────────────────────────────────────
_SEG_RGB = [
    RGBColor(0x4C, 0x72, 0xB0),
    RGBColor(0xDD, 0x8A, 0x52),
    RGBColor(0x55, 0xA8, 0x68),
    RGBColor(0xC4, 0x4E, 0x52),
    RGBColor(0x81, 0x72, 0xB3),
    RGBColor(0x93, 0x7C, 0x60),
    RGBColor(0xDA, 0x8B, 0xC3),
    RGBColor(0x8C, 0x8C, 0x8C),
    RGBColor(0xBC, 0xBD, 0x22),
    RGBColor(0x17, 0xBE, 0xCF),
]
_BOX_RGB = [
    RGBColor(0xC0, 0x50, 0x4D),   # Bottom 2 — red
    RGBColor(0xFF, 0xD9, 0x66),   # Middle   — yellow
    RGBColor(0x44, 0x72, 0xC4),   # Top 2    — blue
]
_DEFAULT_RGB = RGBColor(0x2E, 0x75, 0xB6)

_SCALE_GROUPS = {
    "scale_7": [
        ("Bottom 2 (1-2)", [1, 2]),
        ("Middle (3-5)",   [3, 4, 5]),
        ("Top 2 (6-7)",    [6, 7]),
    ],
    "scale_5": [
        ("Bottom 2 (1-2)", [1, 2]),
        ("Middle (3)",     [3]),
        ("Top 2 (4-5)",    [4, 5]),
    ],
}

CHART_OPTIONS = {
    "categorical": ["100% Stacked bar", "Bar chart", "Horizontal bar", "Pie chart", "Donut chart"],
    "single":      ["100% Stacked bar", "Bar chart", "Horizontal bar", "Pie chart", "Donut chart"],
    "ordinal":     ["100% Stacked bar", "Bar chart", "Horizontal bar", "Line chart"],
    "numeric":     ["Mean", "Histogram", "Box plot", "Violin plot"],
    "scale_7":     ["100% Stacked bar", "Box Stack", "Mean bar", "Bar chart", "Horizontal bar"],
    "scale_5":     ["100% Stacked bar", "Box Stack", "Mean bar", "Bar chart", "Horizontal bar"],
    "multi":       ["Horizontal bar", "Bar chart"],
    "grid":        ["Box Stack", "Mean bar", "Heatmap"],
    "open":        ["Word count bar"],
}


# ── Helpers ───────────────────────────────────────────────────────────────────

def _xml_safe(s) -> str:
    """Strip XML-invalid characters that corrupt the embedded Excel workbook inside PPTX."""
    s = str(s) if s is not None else ""
    # Remove control characters invalid in XML 1.0 (keep \x09 tab handled below)
    s = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', s)
    # Replace newlines, carriage returns, and tabs with a space
    s = re.sub(r'[\r\n\t]', ' ', s)
    # Collapse multiple spaces
    s = re.sub(r'  +', ' ', s)
    return s.strip()


def _f(v) -> float:
    """Convert any numeric to plain Python float (avoids numpy XML issues)."""
    try:
        return float(v)
    except Exception:
        return 0.0


def _short(s, n: int = 40) -> str:
    s = _xml_safe(s)
    return s if len(s) <= n else s[:n - 1] + "…"


def _clean_cats(cats: list) -> list:
    """Ensure categories are plain Python strings, no NaN, XML-safe."""
    return [_xml_safe(c) if (c is not None and str(c) != "nan") else "N/A" for c in cats]


def _value_labels(tile: dict, qnr_questions: list) -> dict:
    labels = {str(k): _xml_safe(v) for k, v in tile.get("value_labels", {}).items()}
    if labels:
        return labels
    raw_opts = tile.get("options", [])
    if raw_opts:
        if isinstance(raw_opts[0], dict):
            for i, opt in enumerate(raw_opts):
                vc = _xml_safe(opt.get("value_coding", "")).strip()
                ol = _xml_safe(opt.get("option_label", "")).strip()
                if "=" in vc:
                    for part in vc.split(","):
                        if "=" in part:
                            k, v = part.strip().split("=", 1)
                            labels.setdefault(k.strip(), v.strip())
                elif vc and ol:
                    labels.setdefault(vc, ol)
                elif ol:
                    labels.setdefault(str(i + 1), ol)
        elif isinstance(raw_opts[0], str):
            for i, opt in enumerate(raw_opts):
                if opt:
                    labels.setdefault(str(i + 1), _xml_safe(opt))
    if labels:
        return labels
    code = tile.get("code", "").upper()
    for q in (qnr_questions or []):
        if q.get("code", "").upper() == code:
            for i, opt in enumerate(q.get("options", [])):
                if isinstance(opt, str) and opt:
                    labels.setdefault(str(i + 1), _xml_safe(opt))
            break
    return labels


def _col_to_label(tile: dict, qnr_questions: list) -> dict:
    all_cols = tile.get("all_cols", [])
    code     = tile.get("code", "").upper()
    opts     = []
    for q in (qnr_questions or []):
        if q.get("code", "").upper() == code:
            opts = q.get("options", [])
            break
    if not opts:
        raw_opts = tile.get("options", [])
        if raw_opts:
            if isinstance(raw_opts[0], dict):
                opts = [o.get("option_label", "") for o in raw_opts]
            elif isinstance(raw_opts[0], str):
                opts = raw_opts
    return {col: _xml_safe(opts[i] if i < len(opts) and opts[i] else col)
            for i, col in enumerate(all_cols)}


# ── Chart data builders ───────────────────────────────────────────────────────

def _data_100pct(series, tile, qnr_questions):
    labels = _value_labels(tile, qnr_questions)
    mapped = series.dropna().map(lambda x: labels.get(str(x), _xml_safe(x)))
    total  = len(mapped)
    if total == 0:
        return None, None
    counts = mapped.value_counts().sort_index()
    if counts.empty:
        return None, None
    data = CategoryChartData()
    data.categories = [_xml_safe(tile.get("code", "Q"))]
    for cat, cnt in counts.items():
        data.add_series(_short(cat), (_f(cnt / total * 100),))
    return data, XL_CHART_TYPE.BAR_STACKED_100


def _data_box_stack(df, tile, qnr_questions):
    col_map    = _col_to_label(tile, qnr_questions)
    cols       = [c for c in tile.get("all_cols", []) if c in df.columns]
    if not cols:
        return None, None
    bar_labels = _clean_cats([col_map.get(c, c) for c in cols])
    var_type   = tile.get("var_type", "scale_7")
    if var_type not in _SCALE_GROUPS:
        # Detect from data
        sample  = df[cols].apply(pd.to_numeric, errors="coerce")
        max_val = sample.max().max()
        var_type = "scale_7" if (pd.notna(max_val) and max_val > 5) else "scale_5"
    groups = _SCALE_GROUPS[var_type]
    data   = CategoryChartData()
    data.categories = bar_labels
    for group_name, rating_vals in groups:
        pcts = []
        for c in cols:
            s     = pd.to_numeric(df[c], errors="coerce").dropna()
            total = len(s)
            pcts.append(_f(s.isin(rating_vals).sum() / total * 100) if total else 0.0)
        data.add_series(_xml_safe(group_name), pcts)
    return data, XL_CHART_TYPE.BAR_STACKED_100


def _data_bar(series, tile, qnr_questions, show_pct, horizontal=False):
    labels = _value_labels(tile, qnr_questions)
    mapped = series.dropna().map(lambda x: labels.get(str(x), _xml_safe(x)))
    total  = len(mapped)
    if total == 0:
        return None, None
    counts = mapped.value_counts().sort_index()
    if counts.empty:
        return None, None
    vals  = [_f(v / total * 100) if show_pct else _f(v) for v in counts.values]
    data  = CategoryChartData()
    data.categories = _clean_cats(list(counts.index))
    data.add_series("% Respondents" if show_pct else "Count", vals)
    ct = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    return data, ct


def _data_multi(df, tile, qnr_questions, show_pct):
    col_map    = _col_to_label(tile, qnr_questions)
    cols       = [c for c in tile.get("all_cols", []) if c in df.columns]
    if not cols:
        return None, None
    total      = len(df)
    raw        = [_f(pd.to_numeric(df[c], errors="coerce").sum()) for c in cols]
    vals       = [_f(v / total * 100) if show_pct else v for v in raw]
    bar_labels = _clean_cats([col_map.get(c, c) for c in cols])
    data       = CategoryChartData()
    data.categories = bar_labels
    data.add_series("% Respondents" if show_pct else "Count", vals)
    return data, XL_CHART_TYPE.BAR_CLUSTERED


def _data_mean_bar(series, tile):
    vals = pd.to_numeric(series, errors="coerce").dropna()
    if len(vals) == 0:
        return None, None
    data = CategoryChartData()
    data.categories = [_xml_safe(tile.get("code", "Mean"))]
    data.add_series("Mean", (_f(vals.mean()),))
    return data, XL_CHART_TYPE.COLUMN_CLUSTERED


def _data_pie(series, tile, qnr_questions):
    labels = _value_labels(tile, qnr_questions)
    mapped = series.dropna().map(lambda x: labels.get(str(x), _xml_safe(x)))
    total  = len(mapped)
    if total == 0:
        return None, None
    counts = mapped.value_counts().sort_index()
    if counts.empty:
        return None, None
    vals  = [_f(v / total * 100) for v in counts.values]
    data  = CategoryChartData()
    data.categories = _clean_cats(list(counts.index))
    data.add_series("", vals)
    return data, XL_CHART_TYPE.PIE


# ── Chart styling ─────────────────────────────────────────────────────────────

def _style_chart(chart, color_mode: str):
    """Apply colours and legend. color_mode: 'single' | 'segments' | 'box'"""
    palette = _BOX_RGB if color_mode == "box" else _SEG_RGB
    n       = len(list(chart.series))

    # Legend
    show_legend = (n > 1)
    try:
        chart.has_legend = show_legend
        if show_legend:
            chart.legend.position         = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
    except Exception:
        pass

    # Series colours
    for i, series in enumerate(chart.series):
        try:
            series.format.fill.solid()
            if n == 1 and color_mode == "single":
                series.format.fill.fore_color.rgb = _DEFAULT_RGB
            else:
                series.format.fill.fore_color.rgb = palette[i % len(palette)]
        except Exception:
            pass


def _add_data_labels(chart, inside: bool = True):
    """Add data labels."""
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(8)
        try:
            dl.position = (XL_LABEL_POSITION.INSIDE_END if inside
                           else XL_LABEL_POSITION.OUTSIDE_END)
        except Exception:
            pass
    except Exception:
        pass


# ── Add one chart block to a slide ────────────────────────────────────────────

def _add_chart_block(slide, tile, df, overrides, qnr_questions,
                     left, top, width, height):
    code     = tile["code"]
    var_type = tile.get("var_type", "categorical")
    g_type   = tile.get("group_type", "single")
    col      = tile.get("dataset_col", "")

    chart_opts = CHART_OPTIONS.get(var_type, ["Bar chart"])
    chosen_ct  = overrides.get(code, chart_opts[0])
    if chosen_ct not in chart_opts:
        chosen_ct = chart_opts[0]
    show_pct = (overrides.get(code + "__pct", "pct") == "pct")

    is_scale     = var_type in ("scale_7", "scale_5")
    is_multi_col = len(tile.get("all_cols", [])) > 1

    # ── Title text box (above chart) ─────────────────────────────────────
    title_top = top - _TITLE_H
    try:
        txb = slide.shapes.add_textbox(left, title_top, width, _TITLE_H)
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = _xml_safe(f"{code} — {tile.get('question', '')[:90]}")
        run.font.size  = Pt(9)
        run.font.bold  = True
        run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
    except Exception:
        pass

    # ── Build chart data ──────────────────────────────────────────────────
    chart_data = None
    pptx_type  = None
    color_mode = "single"

    try:
        if (is_scale and is_multi_col) or g_type == "grid":
            chart_data, pptx_type = _data_box_stack(df, tile, qnr_questions)
            color_mode = "box"

        elif g_type == "multi" and is_multi_col:
            cols   = [c for c in tile.get("all_cols", []) if c in df.columns]
            sample = df[cols].apply(pd.to_numeric, errors="coerce").dropna(how="all")
            if (not sample.empty and pd.notna(sample.max().max())
                    and sample.min().min() >= 1 and sample.max().max() > 2):
                chart_data, pptx_type = _data_box_stack(df, tile, qnr_questions)
                color_mode = "box"
            else:
                chart_data, pptx_type = _data_multi(df, tile, qnr_questions, show_pct)

        elif col and col in df.columns:
            if chosen_ct == "Mean" or var_type == "numeric":
                chart_data, pptx_type = _data_mean_bar(df[col], tile)
            elif chosen_ct == "100% Stacked bar":
                chart_data, pptx_type = _data_100pct(df[col], tile, qnr_questions)
                color_mode = "segments"
            elif chosen_ct in ("Pie chart", "Donut chart"):
                chart_data, pptx_type = _data_pie(df[col], tile, qnr_questions)
                color_mode = "segments"
            elif chosen_ct == "Horizontal bar":
                chart_data, pptx_type = _data_bar(df[col], tile, qnr_questions,
                                                   show_pct, horizontal=True)
            else:
                chart_data, pptx_type = _data_bar(df[col], tile, qnr_questions,
                                                   show_pct, horizontal=False)
    except Exception:
        chart_data = None

    if chart_data is None or pptx_type is None:
        # Fallback: plain text box saying no data
        try:
            fb = slide.shapes.add_textbox(left, top, width, height)
            fb.text_frame.text = _xml_safe(f"{code}: chart could not be generated")
            fb.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xAA, 0x00, 0x00)
        except Exception:
            pass
        return

    # ── Add chart shape ───────────────────────────────────────────────────
    try:
        sp    = slide.shapes.add_chart(pptx_type, left, top, width, height, chart_data)
        chart = sp.chart
        _style_chart(chart, color_mode)
        inside = pptx_type in (XL_CHART_TYPE.BAR_STACKED_100,
                               XL_CHART_TYPE.COLUMN_STACKED_100)
        _add_data_labels(chart, inside=inside)
    except Exception:
        pass


# ── Slide layout helper ───────────────────────────────────────────────────────

def _chart_rects(n: int):
    """
    Return list of (left, top, width, height) for n charts on one slide.
    Layout: up to 3 per row, 2 rows for 4+.
    """
    usable_w = _W - 2 * _MARGIN
    usable_h = _H - _TOP - _MARGIN

    cols = min(n, 3)
    rows = -(-n // cols)   # ceiling division

    col_w = (usable_w - _GAP * (cols - 1)) / cols
    row_h = (usable_h - _GAP * (rows - 1)) / rows

    rects = []
    for i in range(n):
        c = i % cols
        r = i // cols
        left   = _MARGIN + c * (col_w + _GAP)
        top    = _TOP + r * (row_h + _GAP) + _TITLE_H
        height = row_h - _TITLE_H
        rects.append((left, top, col_w, height))
    return rects


def _blank_layout(prs):
    """Return the first blank slide layout, or layout[0] as fallback."""
    for layout in prs.slide_layouts:
        if layout.name.lower() in ("blank", ""):
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


# ── Main entry point ──────────────────────────────────────────────────────────

def build_pptx(visible: list, df, overrides: dict, qnr_questions: list) -> io.BytesIO:
    """
    Build a Presentation with native charts grouped by slide number.
    Tiles without a slide number get one slide each (auto-assigned).
    Returns a seeked BytesIO buffer.
    """
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    blank = _blank_layout(prs)

    # Group tiles by assigned slide number
    slide_groups: dict[int, list] = {}
    for i, tile in enumerate(visible, start=1):
        raw = overrides.get(tile["code"] + "__slide")
        try:
            slide_num = int(float(raw)) if raw not in (None, "") else i
        except (ValueError, TypeError):
            slide_num = i
        slide_groups.setdefault(slide_num, []).append(tile)

    for slide_num in sorted(slide_groups.keys()):
        tiles = slide_groups[slide_num]
        slide = prs.slides.add_slide(blank)
        rects = _chart_rects(len(tiles))
        for tile, (left, top, width, height) in zip(tiles, rects):
            _add_chart_block(slide, tile, df, overrides, qnr_questions,
                             left, top, width, height)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
